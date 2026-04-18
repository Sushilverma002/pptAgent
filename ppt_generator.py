# # ppt_generator.py
# from pptx import Presentation

# def create_ppt(data, output="output.pptx"):
#     prs = Presentation()

#     for slide_data in data["slides"]:
#         slide_layout = prs.slide_layouts[1]
#         slide = prs.slides.add_slide(slide_layout)

#         slide.shapes.title.text = slide_data["title"]
#         slide.placeholders[1].text = "\n".join(slide_data["content"])

#     prs.save(output)
#     return output


# ppt_generator.py

from pptx import Presentation

def create_ppt(data, output="output.pptx", template_path=None):
    
    # Use template if provided
    if template_path:
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    for slide_data in data["slides"]:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # Title
        slide.shapes.title.text = slide_data.get("title", "")

        # Content
        content = slide.placeholders[1]
        content.text = "\n".join(slide_data.get("content", []))

    prs.save(output)
    return output