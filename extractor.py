# extractor.py

from pptx import Presentation

def extract_ppt(path):
    prs = Presentation(path)
    slides = []

    for slide in prs.slides:
        title = ""
        content = []

        if slide.shapes.title:
            title = slide.shapes.title.text

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    content.append(text)

        slides.append({
            "title": title,
            "content": content
        })

    return {"slides": slides}